"""文件读取工具 — 支持图片 base64 编码和文档文本提取。"""

from __future__ import annotations

import base64
import mimetypes
import logging
import shutil
import uuid
from pathlib import Path

logger = logging.getLogger("wokbee")

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"}
DOC_EXTS = {".txt", ".md", ".pdf", ".docx", ".xlsx", ".pptx"}
SUPPORTED_EXTENSIONS = IMAGE_EXTS | DOC_EXTS

_MIME_MAP = {
    ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".png": "image/png", ".gif": "image/gif",
    ".webp": "image/webp", ".bmp": "image/bmp",
}

_ATTACH_ROOT = Path.home() / ".wokbee" / "attachments"


def is_image(path: str) -> bool:
    return Path(path).suffix.lower() in IMAGE_EXTS


def is_document(path: str) -> bool:
    return Path(path).suffix.lower() in DOC_EXTS


def attachments_dir(session_id: str = "tmp") -> Path:
    d = _ATTACH_ROOT / (session_id or "tmp")
    d.mkdir(parents=True, exist_ok=True)
    return d


def save_qimage(image, session_id: str = "tmp", *, suffix: str = ".png") -> str:
    """将 QImage / QPixmap 保存到附件目录，返回路径。"""
    from PySide6.QtGui import QImage, QPixmap

    if isinstance(image, QPixmap):
        image = image.toImage()
    if not isinstance(image, QImage) or image.isNull():
        raise ValueError("无效的图片数据")
    path = attachments_dir(session_id) / f"paste_{uuid.uuid4().hex[:10]}{suffix}"
    if not image.save(str(path), "PNG"):
        raise OSError(f"保存粘贴图片失败: {path}")
    return str(path)


def persist_attachment(src: str, session_id: str = "tmp") -> str:
    """将附件复制到会话目录，保证历史可回看。已在目录内则原样返回。"""
    src_p = Path(src)
    if not src_p.is_file():
        return src
    dest_dir = attachments_dir(session_id)
    try:
        if src_p.resolve().parent == dest_dir.resolve():
            return str(src_p)
    except OSError:
        pass
    dest = dest_dir / f"{uuid.uuid4().hex[:8]}_{src_p.name}"
    shutil.copy2(src_p, dest)
    return str(dest)


def read_image_as_base64(path: str) -> tuple[str, str]:
    """返回 (base64_data, mime_type)。"""
    p = Path(path)
    ext = p.suffix.lower()
    mime = _MIME_MAP.get(ext) or mimetypes.guess_type(str(p))[0] or "image/png"
    data = p.read_bytes()
    return base64.b64encode(data).decode("ascii"), mime


def read_file_as_text(path: str) -> str:
    """根据文件后缀提取纯文本内容。"""
    p = Path(path)
    ext = p.suffix.lower()

    if ext in (".txt", ".md"):
        return _read_plain(p)
    elif ext == ".pdf":
        return _read_pdf(p)
    elif ext == ".docx":
        return _read_docx(p)
    elif ext == ".xlsx":
        return _read_xlsx(p)
    elif ext == ".pptx":
        return _read_pptx(p)
    else:
        return _read_plain(p)


def _read_plain(p: Path) -> str:
    try:
        return p.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return p.read_text(encoding="gbk", errors="replace")


def _read_pdf(p: Path) -> str:
    try:
        import fitz
    except ImportError:
        raise RuntimeError("未安装 PyMuPDF，请执行: pip install PyMuPDF")
    doc = fitz.open(str(p))
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return "\n".join(pages)


def _read_docx(p: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("未安装 python-docx，请执行: pip install python-docx")
    doc = Document(str(p))
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def _read_xlsx(p: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        raise RuntimeError("未安装 openpyxl，请执行: pip install openpyxl")
    wb = load_workbook(str(p), read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def _read_pptx(p: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("未安装 python-pptx，请执行: pip install python-pptx")
    prs = Presentation(str(p))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def build_file_filter() -> str:
    """返回 QFileDialog 使用的文件过滤器字符串。"""
    img_exts = " ".join(f"*{e}" for e in sorted(IMAGE_EXTS))
    doc_exts = " ".join(f"*{e}" for e in sorted(DOC_EXTS))
    all_exts = " ".join(f"*{e}" for e in sorted(SUPPORTED_EXTENSIONS))
    return (
        f"所有支持的文件 ({all_exts});;"
        f"图片文件 ({img_exts});;"
        f"文档文件 ({doc_exts});;"
        "所有文件 (*)"
    )
